import io
import requests
from io import BytesIO
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
import fitz  


def extract_json_array(raw_text):
    """Removes markdown and attempts to safely extract a JSON array."""
    try:
        # Remove markdown formatting if the AI added it
        cleaned = re.sub(r'```json|```', '', raw_text).strip()
        start = cleaned.find('[')
        end = cleaned.rfind(']')
        if start != -1 and end != -1:
            return json.loads(cleaned[start : end + 1])
        return json.loads(cleaned)
    except Exception as e:
        print(f"JSON Parsing Error: {e}\nRaw Text: {raw_text}")
        return []



def extract_text_from_pdf(file_url):
    """Extracts digital text from a PDF."""
    try:
        response = requests.get(file_url, timeout=30)
        response.raise_for_status()
        
        pdf_file = io.BytesIO(response.content)
        reader = PdfReader(pdf_file)
        
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        
        return text.strip()
    except Exception as e:
        print(f"PDF digital extraction error: {e}")
        return ""

def get_pdf_as_images(file_url, max_pages=10):
    """Converts a scanned PDF into a list of PIL Images for Gemini to read natively."""
    try:
        response = requests.get(file_url, timeout=30)
        response.raise_for_status()
        
        # Open PDF from bytes
        doc = fitz.open(stream=response.content, filetype="pdf")
        images = []
        
        # Limit to max_pages to prevent overloading Gemini Quotas
        for page_num in range(min(len(doc), max_pages)):
            page = doc.load_page(page_num)
            # Render page to an image (dpi=150 is good enough for Gemini OCR)
            pix = page.get_pixmap(dpi=150)
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            images.append(img)
            
        return images
    except Exception as e:
        print(f"PyMuPDF rendering error: {e}")
        return []

def extract_text_from_docx(file_url):
    try:
        response = requests.get(file_url, timeout=30)
        response.raise_for_status()
        docx_file = io.BytesIO(response.content)
        doc = Document(docx_file)
        text = "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip())
        return text.strip()
    except Exception as e:
        print(f"DOCX extraction error: {e}")
        return ""

def extract_text_from_pptx(file_url):
    try:
        response = requests.get(file_url, timeout=30)
        response.raise_for_status()
        pptx_file = io.BytesIO(response.content)
        prs = Presentation(pptx_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"PPTX extraction error: {e}")
        return ""
    
def extract_text_from_excel(file_url):
    try:
        response = requests.get(file_url, timeout=30)
        response.raise_for_status()
        excel_file = io.BytesIO(response.content)
        wb = load_workbook(excel_file)
        text = ""
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    text += row_text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Excel extraction error: {e}")
        return ""