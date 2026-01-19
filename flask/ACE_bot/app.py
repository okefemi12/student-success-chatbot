import os
import json
import math
import random
import cohere
from flask import Flask, render_template, request, jsonify
from flask_cors import CORS
import firebase_admin
from firebase_admin import credentials, auth, firestore
from dotenv import load_dotenv
from datetime import datetime, timedelta
from time import sleep
import pytz
from langchain.memory.buffer import ConversationBufferMemory
from langchain_core.chat_history import InMemoryChatMessageHistory
from langchain.schema import HumanMessage, AIMessage
import google.generativeai as genai
from pptx import Presentation
from docx import Document
from openai import OpenAI
import io
import requests
from PyPDF2 import PdfReader
from serpapi.google_search import GoogleSearch
import soundfile as sf
import uuid
import html
import re
import fitz  
import pytesseract
from PIL import Image
import joblib
import numpy as np
import pandas as pd
from sklearn.pipeline import Pipeline
import xgboost as xgb
import cloudinary
import cloudinary.uploader
import requests
from PyPDF2 import PdfReader
from openpyxl import load_workbook
from werkzeug.utils import secure_filename
import warnings
import struct
import wave
from io import BytesIO
import base64
import mimetypes
from key_manager import KeyManager 

warnings.filterwarnings("ignore", category=UserWarning)
warnings.filterwarnings("ignore", category=FutureWarning)
warnings.filterwarnings("ignore", category=DeprecationWarning)


load_dotenv()

# In loop (e.g., PDF chunking), insert a delay:
sleep(random.uniform(1.5, 2.5))  # Prevents hitting quota too fast

# ---------- Flask + Firebase init ----------
app = Flask(__name__, template_folder="templates")
CORS(app)
#firebase authentication
#if not firebase_admin._apps:
    #cred = credentials.Certificate("ace-the-tutor-firebase-adminsdk-fbsvc-737bf9f236.json")
    #firebase_admin.initialize_app(cred)
if not firebase_admin._apps:
    firebase_json = os.getenv("FIREBASE_CREDENTIALS")
    if not firebase_json:
        raise ValueError("Missing FIREBASE_CREDENTIALS environment variable.")
    cred_dict = json.loads(firebase_json)
    cred = credentials.Certificate(cred_dict)
    firebase_admin.initialize_app(cred)


db = firestore.client()

def verify_token_from_request():
    """Verify Firebase ID token from Authorization header (Bearer <token>) or JSON body idToken."""
    auth_header = request.headers.get("Authorization", "")
    token = None
    if auth_header.startswith("Bearer "):
        token = auth_header.split(" ", 1)[1]
    else:
        # fallback: token in JSON body (useful for some POSTs)
        try:
            token = (request.get_json() or {}).get("idToken")
        except Exception:
            token = None

    if not token:
        raise ValueError("Missing ID token")

    decoded = auth.verify_id_token(token)  # fixed
    return decoded


def update_streak(uid):
    """Update user's current and longest streak based on last active date."""
    user_ref = db.collection("users").document(uid)
    user_doc = user_ref.get()
    if not user_doc.exists:
        return

    data = user_doc.to_dict()
    today = datetime.now(pytz.UTC).date()

    last_active_str = data.get("last_active_date")
    current_streak = data.get("streak", 0)
    longest_streak = data.get("longest_streak", 0)

    if last_active_str:
        last_active_date = datetime.strptime(last_active_str, "%Y-%m-%d").date()
        days_diff = (today - last_active_date).days

        if days_diff == 0:
            # already updated today — no change
            return
        elif days_diff == 1:
            current_streak += 1  # continued streak
        else:
            current_streak = 1   # reset streak
    else:
        current_streak = 1  # first time

    longest_streak = max(longest_streak, current_streak)

    user_ref.update({
        "streak": current_streak,
        "longest_streak": longest_streak,
        "last_active_date": today.strftime("%Y-%m-%d")
    })

# ---------- JSON Serializer ----------
def clean_value(value):
    """Convert Firestore/special values into JSON-safe formats."""
    if isinstance(value, datetime):
        return value.isoformat()
    if value == firestore.SERVER_TIMESTAMP:
        return None  # placeholder until Firestore resolves
    return value

def safe_clean_dict(d):
    """Recursively clean dicts/lists for JSON serialization."""
    if isinstance(d, dict):
        return {k: safe_clean_dict(v) for k, v in d.items()}
    elif isinstance(d, list):
        return [safe_clean_dict(v) for v in d]
    else:
        return clean_value(d)
    
groq_client = OpenAI(
                   api_key=os.getenv("GROQ_API_KEY"),
                   base_url="https://api.groq.com/openai/v1"
                )
    

#langchain memory  (Chatbot's memory)
conversation_memories = {}

def get_memory(session_id):
    if session_id not in conversation_memories:

        history = InMemoryChatMessageHistory()
        conversation_memories[session_id] = ConversationBufferMemory(
            return_messages=True,
            memory_key="chat_history",
            chat_memory=history
        )
    return conversation_memories[session_id]

#For clean responses 

def clean_response(text):
    # Unescape HTML entities
    text = html.unescape(text)
    # Replace escaped newline and tab characters
    text = text.replace('\\n', '\n').replace('\\t', '\t')
    # Replace escaped quotes
    text = text.replace('\\"', '"').replace("\\'", "'")
    # Collapse multiple whitespaces
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)  
    text = re.sub(r'^\* ', '• ', text, flags=re.MULTILINE) 

    # Convert markdown to HTML
    text = re.sub(r'\*\*(.*?)\*\*', r'<strong>\1</strong>', text)
    text = re.sub(r'\*(.*?)\*', r'<em>\1</em>', text)
    
    # Clean spacing
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r'\n\n\n+', '\n\n', text)
    return text.strip()


UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

#api keys
serper_api_key = os.getenv("SERPER_API_KEY")
genai.configure(api_key=os.getenv("GEMINI_API_KEY_1"))
model = genai.GenerativeModel(model_name="gemini-2.5-flash")
genai.configure(api_key=os.getenv("RECOMMEND_BACKUP1"))
rec_backup1 = genai.GenerativeModel(model_name="gemini-2.5-flash")
genai.configure(api_key=os.getenv("RECOMMEND_BACKUP2"))
rec_backup2 = genai.GenerativeModel(model_name="gemini-2.5-flash")
genai.configure(api_key=os.getenv("GEMINI_API_KEY_2"))
model_pdf = genai.GenerativeModel(model_name="gemini-2.5-flash")
genai.configure(api_key=os.getenv("GEMINI_API_KEY_3"))
model_cards= genai.GenerativeModel(model_name="gemini-2.5-flash")
genai.configure(api_key=os.getenv("GEMINI_API_KEY_4"))
reminder_model = genai.GenerativeModel("gemini-2.5-flash")
TTS_API_KEYS = [
    os.environ.get("ACE_VOICE"),              # Primary
    os.environ.get("ACE_VOICE_BACKUP_1"),     # Backup 1
    os.environ.get("ACE_VOICE_BACKUP_2"),     # Backup 2
    os.environ.get("ACE_VOICE_BACKUP_3"),     # Backup 3
]
FINAL_BACKUP = [
    os.environ.get("BACKUP_1"),             
    os.environ.get("BACKUP_2"),     
    os.environ.get("BACKUP_3"),    
    os.environ.get("BACKUP_4"),
    os.environ.get("BACKUP_5"),
    os.environ.get("BACKUP_6"),
]
backup_manager = KeyManager(FINAL_BACKUP, model_name="gemini-2.5-flash")


YOUTUBE_API_KEY = os.getenv("YOUTUBE_API_KEY")
cloudinary.config(
    cloud_name=os.getenv("CLOUD_NAME"),
    api_key=os.getenv("API_KEY"),
    api_secret=os.getenv("API_SECRET")
)

def rec_backup_response(prompt):
    try:
        return rec_backup1.generate_content(prompt)
    except Exception:
        try: 
            return rec_backup2.generate_content(prompt)
        except Exception as e: 
            return f"Error: {e}"
#Backup models 
def generate_backup_response(prompt):
    try:
        # Try Gemini
        return model.generate_content(prompt).text
    except Exception:
        try:
            # Try Cohere
            co = cohere.Client(os.getenv("COHERE_API_KEY"))
            response = co.chat(model="command-r-plus-08-2024", message=prompt)
            return response.text.strip()
        except Exception:
            try:
                client = OpenAI(
                   api_key=os.getenv("GROQ_API_KEY"),
                   base_url="https://api.groq.com/openai/v1"
                )

                groq_response = client.chat.completions.create(
                    model="llama-3.1-8b-instant",
                    messages=[{"role": "user", "content": prompt}]
                )

                return groq_response.choices[0].message.content
            except Exception as e:
                return f"All model attempts failed: {e}"
#search reasources via internet 
def search_web(query):
    headers = {
        "X-API-KEY": serper_api_key,  # replace this
        "Content-Type": "application/json"
    }

    data = {
        "q": query
    }

    response = requests.post("https://google.serper.dev/search", headers=headers, json=data)

    if response.status_code != 200:
        return [{"title": "Error", "link": f"Status code: {response.status_code}"}]

    results = response.json()

    links = []
    for result in results.get("organic", []):
        links.append({
            "title": result.get("title"),
            "link": result.get("link")
        })

    return links



# --- BACKUP FUNCTION ---
def final_backup_response(prompt_or_content):
    """
    Tries to generate response using the backup key rotation.
    Accepts plain text OR [prompt, image] list.
    """
    try:
        # Calls the KeyManager to rotate through your 6 FINAL_BACKUP keys
        response = backup_manager.generate_content(prompt_or_content)
        return response.text
    except Exception as e:
        print(f"All 6 Backup Keys failed: {e}")
        return "Error: Service temporarily unavailable. Please try again later."
# ---------- Helper: Image OCR ----------

def extract_text_from_image(image_url):
    try:
        # 1. Download the image bytes from Cloudinary
        response = requests.get(image_url)
        response.raise_for_status() # Check for download errors
        
        # 2. Convert bytes to an Image object
        img = Image.open(BytesIO(response.content))
        
        # 3. Extract text
        text = pytesseract.image_to_string(img)
        return text.strip()
    except Exception as e:
        print(f"OCR Error: {e}")
        return ""
    

def extract_text_from_pdf(file_url):
    """Extract text from PDF URL without saving to disk"""
    try:
        response = requests.get(file_url, timeout=30)
        response.raise_for_status()
        
        # Read PDF directly from bytes in memory
        pdf_file = io.BytesIO(response.content)
        reader = PdfReader(pdf_file)
        
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        
        return text.strip()
        
    except Exception as e:
        print(f"PDF extraction error: {e}")
        return ""


def extract_text_from_docx(file_url):
    """Extract text from DOCX URL without saving to disk"""
    try:
        response = requests.get(file_url, timeout=30)
        response.raise_for_status()
        
        # Read DOCX directly from bytes in memory
        docx_file = io.BytesIO(response.content)
        doc = Document(docx_file)
        
        text = "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip())
        return text.strip()
        
    except Exception as e:
        print(f"DOCX extraction error: {e}")
        return ""


def extract_text_from_pptx(file_url):
    """Extract text from PPTX URL without saving to disk"""
    try:
        response = requests.get(file_url, timeout=30)
        response.raise_for_status()
        
        # Read PPTX directly from bytes in memory
        pptx_file = io.BytesIO(response.content)
        prs = Presentation(pptx_file)
        
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return text.strip()
        
    except Exception as e:
        print(f"PPTX extraction error: {e}")
        return ""


def extract_text_from_excel(file_url):
    """Extract text from Excel URL without saving to disk"""
    try:
        response = requests.get(file_url, timeout=30)
        response.raise_for_status()
        
        # Read Excel directly from bytes in memory
        excel_file = io.BytesIO(response.content)
        wb = load_workbook(excel_file)
        
        text = ""
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    text += row_text + "\n"
        
        return text.strip()
        
    except Exception as e:
        print(f"Excel extraction error: {e}")
        return ""


TTS_API_KEYS = [key for key in TTS_API_KEYS if key]

# Track current TTS key index
current_tts_key_index = 0



# DD THIS FUNCTION FIRST (before generate_audio_with_retry)
def convert_to_wav(pcm_data, mime_type):
    """
    Convert raw PCM audio data to WAV format.
    
    Args:
        pcm_data: Raw PCM audio bytes
        mime_type: MIME type string (e.g., "audio/L16;rate=24000")
    
    Returns:
        bytes: WAV file data
    """
    try:
        # Parse sample rate from mime_type (e.g., "audio/L16;rate=24000")
        sample_rate = 24000  # default
        if "rate=" in mime_type:
            rate_str = mime_type.split("rate=")[1].split(";")[0]
            sample_rate = int(rate_str)
        
        # PCM parameters
        num_channels = 1  # mono
        sample_width = 2  # 16-bit = 2 bytes
        
        # Create WAV file in memory
        wav_buffer = io.BytesIO()
        
        with wave.open(wav_buffer, 'wb') as wav_file:
            wav_file.setnchannels(num_channels)
            wav_file.setsampwidth(sample_width)
            wav_file.setframerate(sample_rate)
            wav_file.writeframes(pcm_data)
        
        # Get WAV data
        wav_buffer.seek(0)
        return wav_buffer.read()
        
    except Exception as e:
        print(f"⚠️ WAV conversion error: {e}")
        # Return original data if conversion fails
        return pcm_data


# NOW define generate_audio_with_retry (after convert_to_wav)
def generate_audio_with_retry(text, max_retries=None):
    """
    Generate audio with automatic API key rotation on quota OR network errors.
    """
    global current_tts_key_index
    
    if max_retries is None:
        max_retries = len(TTS_API_KEYS)
    
    last_error = None
    
    for attempt in range(max_retries):
        try:
            from google import genai
            from google.genai import types
            
            # Get current TTS API key
            current_key = TTS_API_KEYS[current_tts_key_index]
            print(f"🔑 Attempting with TTS API key #{current_tts_key_index + 1}")
            
            tts_client = genai.Client(api_key=current_key)
            
            tts_contents = [
                types.Content(
                    role="user",
                    parts=[types.Part.from_text(text=text)]
                )
            ]
            
            tts_config = types.GenerateContentConfig(
                response_modalities=["audio"],
                speech_config=types.SpeechConfig(
                    voice_config=types.VoiceConfig(
                        prebuilt_voice_config=types.PrebuiltVoiceConfig(
                            voice_name="Sadachbia"
                        )
                    )
                )
            )
            
            audio_chunks = []
            audio_mime_type = None
            
            # Generate stream
            for chunk in tts_client.models.generate_content_stream(
                model="gemini-2.5-flash", # Ensure model name is correct (remove -preview-tts if needed)
                contents=tts_contents,
                config=tts_config
            ):
                if (chunk.candidates and len(chunk.candidates) > 0 and 
                    chunk.candidates[0].content and chunk.candidates[0].content.parts):
                    
                    part = chunk.candidates[0].content.parts[0]
                    
                    if hasattr(part, 'inline_data') and part.inline_data:
                        if hasattr(part.inline_data, 'data') and part.inline_data.data:
                            audio_chunks.append(part.inline_data.data)
                            
                            if audio_mime_type is None and hasattr(part.inline_data, 'mime_type'):
                                audio_mime_type = part.inline_data.mime_type
            
            if not audio_chunks:
                raise Exception("No audio chunks received (Stream Empty)")
            
            combined_audio = b"".join(audio_chunks)
            
            # Convert to WAV if needed (ensure convert_to_wav is defined in your code)
            if audio_mime_type and "audio/L" in audio_mime_type:
                wav_data = convert_to_wav(combined_audio, audio_mime_type)
                audio_base64 = base64.b64encode(wav_data).decode("utf-8")
                mime_type = "audio/wav"
            else:
                audio_base64 = base64.b64encode(combined_audio).decode("utf-8")
                mime_type = audio_mime_type or "audio/wav"
            
            print(f"TTS success with API key #{current_tts_key_index + 1}")
            return audio_base64, mime_type
                
        except Exception as e:
            error_msg = str(e).lower()
            last_error = e
            
            # Check for Quota limits OR Network Connection errors
            is_quota_error = any(k in error_msg for k in ['quota', 'rate limit', 'resource exhausted', '429'])
            is_network_error = any(k in error_msg for k in ['getaddrinfo', 'connecterror', 'socket', 'errno 11002', '11001'])
            
            if is_quota_error:
                print(f"API key #{current_tts_key_index + 1} exhausted. Switching...")
                current_tts_key_index = (current_tts_key_index + 1) % len(TTS_API_KEYS)
                
            elif is_network_error:
                print(f"Network Error (DNS/Connection) with key #{current_tts_key_index + 1}: {e}")
                print("Waiting 3 seconds for internet to stabilize...")
                time.sleep(3) # Wait before retrying
                # We do NOT switch keys for network errors, we just retry the loop (or you can switch if you prefer)
                # If you want to switch keys anyway:
                # current_tts_key_index = (current_tts_key_index + 1) % len(TTS_API_KEYS)
                
            else:
                # Actual code error (not network/quota), so we crash
                print(f"❌ Unrecoverable error with TTS API key #{current_tts_key_index + 1}: {e}")
                raise e
    
    raise Exception(f"TTS generation failed after {max_retries} attempts. Last error: {last_error}")

def parse_study_plan(raw_text):
    study_plan = []
    for line in raw_text.splitlines():
        match = re.match(r".*Day\s*(\d+)[\:\-]?\s*(.+)", line, re.IGNORECASE)
        if match:
            topics = [t.strip() for t in re.split(r",|;", match.group(2)) if t.strip()]
            study_plan.append({"day": int(match.group(1)), "topics": topics})
    return study_plan
# Load the saved model once at startup
pipeline = joblib.load("model/student_prediction_model.pkl")


@app.route("/")
def home():
    return jsonify({
        "status": "OK",
        "message": "Ace The Tutor backend API is live 🚀"
    })


# ---------- Index ----------
@app.route("/test")
def index():
    firebase_config = {
        "apiKey": os.getenv("FIREBASE_API_KEY"),
        "authDomain": os.getenv("FIREBASE_AUTH_DOMAIN"),
        "projectId": os.getenv("FIREBASE_PROJECT_ID"),
        "storageBucket": os.getenv("FIREBASE_STORAGE_BUCKET"),
        "messagingSenderId": os.getenv("FIREBASE_MESSAGING_SENDER_ID"),
        "appId": os.getenv("FIREBASE_APP_ID"),
    }
    return render_template("index.html", firebase_config=firebase_config)




# ---------- Register ----------
@app.route("/register", methods=["POST"])
def register_user():
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]
        email = decoded.get("email")

        data = request.get_json() or {}
        name = data.get("name") or decoded.get("name") or ""
        date_of_birth = data.get("date_of_birth")
        gender = data.get("gender")
        phone_number = data.get("phone_number") or data.get("phone")

  
        subject = data.get("subject") or []
        if isinstance(subject, str):
            subject = [s.strip() for s in subject.split(",") if s.strip()]

        course_of_study = data.get("course_of_study")
        country = data.get("country")
        school_type = data.get("school_type")
        school_name = data.get("school_name")
        degree = data.get("degree")

        if not name or not date_of_birth or not gender or not phone_number:
            return jsonify({"error": "Required fields: name, date_of_birth, gender, phone_number"}), 400

        user_doc_ref = db.collection("users").document(uid)
        user_doc = user_doc_ref.get()

        if not user_doc.exists:  
            # First time registration → create with defaults
            user_data = {
                "user_id": uid,
                "name": name,
                "email": email,
                "date_of_birth": date_of_birth,
                "gender": gender,
                "phone_number": phone_number,
                "subject": subject,
                "school_type": school_type,
                "school_name": school_name,
                "country": country,
                "degree": degree,
                "course_of_study": course_of_study,
                "created_at": firestore.SERVER_TIMESTAMP,
                "provider": decoded.get("firebase", {}).get("sign_in_provider"),

                # Initialize log activity fields
                "AttendanceDays": [],
                "Absences": 0,
                "streak": 0,
                "attendance_percentage": 0,
                "study_hours_per_week": 0,
                "last_weekly_reset": datetime.now(pytz.UTC).strftime("%Y-%m-%d"),
                "assignment_completed": 0,
                "Tutoring": 0,
                "sleep_hours_per_day": 8,
                "participation_level": 0
            }
            user_doc_ref.set(user_data)
        else:
            # If user already exists, only update profile info
            updates = {
                "name": name,
                "date_of_birth": date_of_birth,
                "gender": gender,
                "phone_number": phone_number,
                "subject": subject,
                "course_of_study": course_of_study,
                "country": country,
                "school_type": school_type,
                "school_name": school_name,
                "degree": degree,
            }
            updates = {k: v for k, v in updates.items() if v is not None}
            user_doc_ref.update(updates)
            user_data = {**user_doc.to_dict(), **updates}

        returned = user_data.copy()
        returned.pop("provider", None)

        return jsonify({"ok": True, "uid": uid, "profile": safe_clean_dict(returned)})
    except Exception as e:
        return jsonify({"error": str(e)}), 400



# ---------- Profile ----------
@app.route("/profile", methods=["GET"])
def get_profile():
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        doc = db.collection("users").document(uid).get()
        if not doc.exists:
            return jsonify({"error": "User profile not found"}), 404

        profile = doc.to_dict()
        profile.pop("provider", None)

        return jsonify({"ok": True, "profile": safe_clean_dict(profile)})
    except Exception as e:
        return jsonify({"error": str(e)}), 401


# ---------- Update Profile ----------
@app.route("/update_profile", methods=["POST"])
def update_profile():
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]
        data = request.get_json() or {}

        
        allowed = [
            "name", "date_of_birth", "gender", "phone_number", "subject",
            "course_of_study", "country", "school_type", "school_name", "degree"
        ]
        updates = {}

        for k in allowed:
            if k in data:
                if k == "subject":
                    # Always store subject as list
                    subject = data["subject"]
                    if isinstance(subject, str):
                        subject = [s.strip() for s in subject.split(",") if s.strip()]
                    updates["subject"] = subject
                else:
                    updates[k] = data[k]

        if not updates:
            return jsonify({"error": "No updatable fields provided"}), 400

        db.collection("users").document(uid).update(updates)

        doc = db.collection("users").document(uid).get()
        profile = doc.to_dict() or {}
        profile.pop("provider", None)

        return jsonify({"ok": True, "profile": safe_clean_dict(profile)})
    except Exception as e:
        return jsonify({"error": str(e)}), 400

# ---------- Delete Account ----------
@app.route("/delete-account", methods=["DELETE"])
def delete_account():
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        db.collection("users").document(uid).delete()
        auth.delete_user(uid)

        return jsonify({"ok": True, "message": "Account deleted successfully"})
    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ---------- Log Activity ----------
@app.route("/log_activity", methods=["POST"])
def log_activity():
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        data = request.get_json(silent=True) or {}
        session_hours = data.get("session_hours")
        if session_hours is None:
            return jsonify({"error": "Missing session_hours"}), 400

        try:
            session_hours = float(session_hours)
        except Exception:
            return jsonify({"error": "session_hours must be a number"}), 400

        if session_hours <= 0:
            return jsonify({"error": "session_hours must be > 0"}), 400

        today_dt = datetime.now(pytz.UTC)
        today_date = today_dt.date()
        today_str = today_date.strftime("%Y-%m-%d")

        user_doc_ref = db.collection("users").document(uid)
        user_doc = user_doc_ref.get()
        if not user_doc.exists:
            return jsonify({"error": "User not found"}), 404

        user_data = user_doc.to_dict()
        attendance_days = user_data.get("AttendanceDays", [])

        # --- Update today's record ---
        today_record = next((d for d in attendance_days if d.get("date") == today_str), None)
        if today_record:
            today_record["hours_used"] = today_record.get("hours_used", 0) + session_hours
            # Mark attended if studied 30+ minutes
            today_record["attended"] = today_record["hours_used"] >= 0.5
        else:
            today_record = {
                "date": today_str,
                "hours_used": session_hours,
                "attended": session_hours >= 0.5 
            }
            attendance_days.append(today_record)

        # --- Attendance metrics ---
        parsed_dates = [datetime.strptime(d["date"], "%Y-%m-%d").date() for d in attendance_days]
        start_date = min(parsed_dates) if attendance_days else today_date
        total_days = (today_date - start_date).days + 1
        attended_days_count = sum(1 for d in attendance_days if d.get("attended"))
        absences = max(0, total_days - attended_days_count)
        attendance_percentage = round((attended_days_count / total_days) * 100, 2) if total_days > 0 else 0

        # --- Weekly reset ---
        last_reset_str = user_data.get("last_weekly_reset")
        last_reset_date = datetime.strptime(last_reset_str, "%Y-%m-%d").date() if last_reset_str else today_date
        if (today_date - last_reset_date).days >= 7:
            study_hours_per_week = session_hours
            last_reset_date = today_date
        else:
            one_week_ago = today_date - timedelta(days=7)
            study_hours_per_week = sum(
                d.get("hours_used", 0)
                for d in attendance_days
                if datetime.strptime(d["date"], "%Y-%m-%d").date() >= one_week_ago
            )

        # --- Fetch quiz performance for assignment completion ---
        quiz_scores_ref = db.collection("users").document(uid).collection("quiz_scores")
        quiz_docs = quiz_scores_ref.stream()
        percentages = [doc.to_dict().get("percentage", 0) for doc in quiz_docs]
        assignments_completed = round(sum(percentages) / (len(percentages) * 100), 3) if percentages else 0.0

        # --- Tutoring recommendation ---
        Tutoring = 1 if study_hours_per_week < 5 or attendance_percentage < 70 else 0

        # --- Participation Level (binary for model) ---
        participation_level = 1 if attendance_percentage >= 50 else 0

        # --- Sleep (default fallback) ---
        sleep_hours_per_day = user_data.get("sleep_hours_per_day", 8)

        # --- Update Firestore ---
        user_doc_ref.update({
            "AttendanceDays": attendance_days,
            "Absences": absences,
            "attendance_percentage": attendance_percentage,
            "study_hours_per_week": study_hours_per_week,
            "last_weekly_reset": last_reset_date.strftime("%Y-%m-%d"),
            "assignment_completed": assignments_completed,
            "participation_level": participation_level,
            "Tutoring": Tutoring
        })

        return jsonify({
            "ok": True,
            "updated_metrics": safe_clean_dict({
                "attendance_percentage": attendance_percentage,
                "study_hours_per_week": study_hours_per_week,
                "assignment_completed": assignments_completed,
                "participation_level": participation_level,
                "Tutoring": Tutoring
            })
        })

    except Exception as e:
        return jsonify({"error": str(e)}), 400



@app.route("/create_chat_session", methods=["POST"])
def create_chat_session():
    try:
        decoded = verify_token_from_request()  
        uid = decoded["uid"]

        # Generate unique session ID
        session_id = str(uuid.uuid4())

        # Save session metadata in Firestore
        session_ref = (
            db.collection("users")
            .document(uid)
            .collection("chat_sessions")
            .document(session_id)
        )
        session_ref.set({
            "created_at": firestore.SERVER_TIMESTAMP,
            "title": None  # will be auto-set on first user message
        })

        return jsonify({"ok": True, "session_id": session_id})
    except Exception as e:
        return jsonify({"error": str(e)}), 400



@app.route("/chat/<session_id>", methods=["POST"])
def chat(session_id):
    try:
        decoded = verify_token_from_request()
        uid = decoded['uid']
        update_streak(uid)

        data = request.get_json() or {}
        user_message = data.get("message", "").strip()
        if not user_message:
            return jsonify({"error": "No message provided"}), 400

        # 🔹 Check if session exists + auto-set title if missing
        session_ref = (
            db.collection("users")
            .document(uid)
            .collection("chat_sessions")
            .document(session_id)
        )
        session_doc = session_ref.get()
        if session_doc.exists:
            session_data = session_doc.to_dict()
            if not session_data.get("title"):
                title = user_message[:50] + ("..." if len(user_message) > 50 else "")
                session_ref.update({"title": title})
        else:
            return jsonify({"error": "Invalid session_id"}), 404

        # 🔹 Load LangChain memory
        memory = get_memory(session_id)
        if not memory.chat_memory.messages:
            docs = (
                session_ref.collection("messages")
                .order_by("timestamp", direction=firestore.Query.ASCENDING)
                .limit(20)
                .stream()
            )
            for d in docs:
                m = d.to_dict()
                if m["role"] == "user":
                    memory.chat_memory.add_user_message(m["content"])
                elif m["role"] == "assistant":
                    memory.chat_memory.add_ai_message(m["content"])

        # 🔹 Build conversation context
        context = ""
        if memory.chat_memory.messages:
            context = "\n\nConversation History:\n"
            for msg in memory.chat_memory.messages[-10:]:
                if isinstance(msg, HumanMessage):
                    context += f"User: {msg.content}\n"
                elif isinstance(msg, AIMessage):
                    context += f"Assistant: {msg.content}\n"

        # 🔹 Fetch user profile
        profile_doc = db.collection("users").document(uid).get()
        profile_data = profile_doc.to_dict() if profile_doc.exists else {}

        profile_context = ""
        if not memory.chat_memory.messages:  # Only for first user input
            profile_doc = db.collection("users").document(uid).get()
            profile_data = profile_doc.to_dict() if profile_doc.exists else {}
            if profile_data:
                profile_context = (
                    f"User profile:\n"
                    f"- Name: {profile_data.get('name', 'N/A')}\n"
                    f"- Course of Study: {profile_data.get('course_of_study', 'N/A')}\n"
                    f"- Subject: {profile_data.get('subject', 'N/A')}\n"
                    f"- Degree: {profile_data.get('degree', 'N/A')}\n"
                    f"- school_type: {profile_data.get('school_type', 'N/A')}\n"
                    f"- school_name: {profile_data.get('school_name', 'N/A')}\n"
                    f"- country: {profile_data.get('country', 'N/A')}\n"
                )
        

        # 🔹 Build Gemini prompt
        gemini_prompt = (
            "You are FELIX, a helpful intelligent AI study tutor and assistant. Your task is to assist with studying. "
            "Your goal is to help students deeply understand and apply academic concepts, not just recall them.\n\n"
    
            "Follow these steps carefully for each question:\n"
            "• Start by explaining the key concept in simple, clear language.\n"
            "• If the question involves a calculation, proof, or problem, show the full worked solution step by step.\n"
            "• After solving, explain what the result means and give a quick real-life or practical connection.\n"
            "• If the user only asks for understanding (not solving), focus on clear explanations and short examples.\n"
            "• Keep your tone encouraging, like a supportive tutor.\n\n"
    
    
            "Guidelines:\n"
            "• Focus strictly on academic and study-related questions — avoid casual conversation.\n"
            "• If a question is unclear or off-topic, politely ask for clarification.\n"
            "• When the user provides complex questions, break them into simpler parts before answering.\n"
            "• Only greet the user with their name in the very first message of a new session.\n"
            "• Don't mention their name or greet again after the your first response.\n"
            "• Format your response with clear paragraphs and bullet points using the • symbol.\n"
            "• Avoid markdown or special characters.\n"
            "• Use previous context to maintain continuity.\n\n"
            "• Use plain ASCII math notation (sqrt, ^, /, lim_{x->0}, etc.).\n"
            
    
            f"{profile_context}\n"
            f"{context}\n\n"
            f"Current question: {user_message}\n\n"
    
            "Now think carefully and respond as FELIX — explain the concept briefly, then solve step by step if needed, "
            "and finish with a clear final answer and one short follow-up question."
            "If a topic is complex, recommend that the user reviews additional resources.\n"
       )


        try:
            gemini_response = model.generate_content(gemini_prompt)
            raw_answer = gemini_response.text.strip()
            answer = clean_response(raw_answer)
        except Exception as e:
            print(f"Gemini failed: {e}")
            raw_answer = generate_backup_response(gemini_prompt)
            answer = clean_response(raw_answer)

        # 🔹 Optional: add web search links
        if any(keyword in user_message.lower() for keyword in ["how to", "explain", "definition", "what is", "research", "video", "reference"]):
            web_links = search_web(user_message)
            if web_links:
                limited_links = web_links[:3]
                links_text = "\n\nUseful Links:\n"
                for link in limited_links:
                    title = link.get("title", "View Resource")
                    url = link.get("link", "")
                    links_text += f"• {title}: {url}\n"
                answer += links_text
                answer = clean_response(answer)

      
        memory.chat_memory.add_user_message(user_message)
        memory.chat_memory.add_ai_message(answer)

        
        chat_ref = session_ref.collection("messages")
        chat_ref.add({
            "role": "user",
            "content": user_message,
            "timestamp": firestore.SERVER_TIMESTAMP
        })
        chat_ref.add({
            "role": "assistant",
            "content": answer,
            "timestamp": firestore.SERVER_TIMESTAMP
        })

        return jsonify({
            "response": answer,
            "session_id": session_id
        })

    except Exception as e:
        return jsonify({"error": str(e)}), 500
@app.route("/chat/<session_id>/audio", methods=["POST"])
def generate_audio(session_id):
    """
    Generate audio for a specific message in the chat session.
    Expects: { "message_id": "...", "text": "..." } or just { "text": "..." }
    """
    try:
        # Verify user authentication
        decoded = verify_token_from_request()
        uid = decoded['uid']
        
        # Verify session belongs to user
        session_ref = (
            db.collection("users")
            .document(uid)
            .collection("chat_sessions")
            .document(session_id)
        )
        session_doc = session_ref.get()
        if not session_doc.exists:
            return jsonify({"error": "Invalid session_id"}), 404
        
        # Get text to convert to audio
        data = request.get_json() or {}
        text = data.get("text", "").strip()
        message_id = data.get("message_id")  # Optional: to update specific message
        
        if not text:
            return jsonify({"error": "No text provided"}), 400
        
        # Generate audio with automatic retry/fallback
        try:
            audio_base64, mime_type = generate_audio_with_retry(text)
            
        except Exception as e:
            print(f"TTS generation failed after all retries: {str(e)}")
            import traceback
            traceback.print_exc()
            return jsonify({"error": f"TTS generation failed: {str(e)}"}), 500
        
        # Optional: Update Firestore message with audio data if message_id provided
        if message_id and audio_base64:
            try:
                message_ref = session_ref.collection("messages").document(message_id)
                message_ref.update({
                    "audio": audio_base64,
                    "mime_type": mime_type,
                    "audio_generated_at": firestore.SERVER_TIMESTAMP
                })
            except Exception as e:
                print(f" Failed to update message with audio: {e}")
        
        # Return audio data
        return jsonify({
            "audio": audio_base64,
            "mime_type": mime_type,
            "message_id": message_id
        })

    except Exception as e:
        print(f" Error in audio generation route: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route('/chat/clear-memory', methods=['POST'])
def clear_memory():
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        data = request.get_json() or {}
        session_id = data.get("session_id")
        if not session_id:
            return jsonify({"error": "Missing session_id"}), 400

        if session_id in conversation_memories:
            conversation_memories[session_id].clear()

        chat_ref = (
            db.collection("users")
            .document(uid)
            .collection("chat_sessions")
            .document(session_id)
            .collection("messages")
        )
        docs = chat_ref.stream()
        for d in docs:
            d.reference.delete()

        return jsonify({"ok": True, "message": f"Memory cleared for session {session_id}"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/chat/get-history', methods=['POST'])
def get_history():
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        data = request.get_json() or {}
        session_id = data.get("session_id")
        if not session_id:
            return jsonify({"error": "Missing session_id"}), 400

        # 🔹 Fetch session metadata
        session_ref = (
            db.collection("users")
            .document(uid)
            .collection("chat_sessions")
            .document(session_id)
        )
        session_doc = session_ref.get()
        if not session_doc.exists:
            return jsonify({"error": "Invalid session_id"}), 404
        session_data = session_doc.to_dict()

        # 🔹 Fetch messages
        chat_ref = session_ref.collection("messages")
        docs = chat_ref.order_by("timestamp", direction=firestore.Query.ASCENDING).stream()

        messages = []
        for d in docs:
            m = d.to_dict()
            messages.append({
                "role": m.get("role"),
                "content": m.get("content"),
                "timestamp": m.get("timestamp").isoformat() if m.get("timestamp") else None
            })

        return jsonify({
            "session_id": session_id,
            "title": session_data.get("title"),
            "messages": messages,
            "total_messages": len(messages)
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/chat_audio/<session_id>", methods=["POST"])
def chat_audio(session_id):
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]
        update_streak(uid)

        # Check audio
        if "audio" not in request.files:
            return jsonify({"error": "No audio file uploaded"}), 400

        audio_file = request.files["audio"]

        # Save temporarily to memory or file system
        audio_bytes = audio_file.read()

        # Step 1: Transcribe with Groq Whisper
        transcription = groq_client.audio.transcriptions.create(
            model="whisper-large-v3",
            file=("audio.wav", audio_bytes)
        )
        user_message = transcription.text.strip()

        # Step 2: Save user message to Firestore
        message_id = str(uuid.uuid4())
        message_doc = {
            "role": "user",
            "content": user_message,
            "created_at": firestore.SERVER_TIMESTAMP,
        }
        db.collection("users").document(uid).collection("chat_sessions") \
            .document(session_id).collection("messages") \
            .document(message_id).set(message_doc)

        # Step 3: Generate chatbot response
        gemini_prompt = f"User said (from audio): {user_message}"
        try:
            gemini_response = model.generate_content(gemini_prompt)
            raw_answer = gemini_response.text.strip()
            answer = clean_response(raw_answer)
        except Exception as e:
            print(f"Gemini failed: {e}")
            raw_answer = generate_backup_response(gemini_prompt)
            answer = clean_response(raw_answer)

        # Step 4: Save bot response to Firestore
        bot_message_id = str(uuid.uuid4())
        bot_message_doc = {
            "role": "assistant",
            "content": answer,
            "created_at": firestore.SERVER_TIMESTAMP,
        }
        db.collection("users").document(uid).collection("chat_sessions") \
            .document(session_id).collection("messages") \
            .document(bot_message_id).set(bot_message_doc)

        return jsonify({
            "ok": True,
            "transcription": user_message,
            "answer": answer
        })

    except Exception as e:
        return jsonify({"error": str(e)}), 500

    

    
 # Make sure this is imported

@app.route('/media_summary', methods=['POST'])
def media_summary():
    try:
        # --- 1. Auth & Setup ---
        decoded = verify_token_from_request()
        uid = decoded["uid"]
        update_streak(uid)

        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400

        file = request.files['file']
        filename = file.filename.lower()

        # Determine Mime Type (Crucial for passing images to Gemini)
        mime_type, _ = mimetypes.guess_type(filename)

        # --- 2. Upload to Cloudinary ---
        upload_result = cloudinary.uploader.upload(
            file,
            upload_preset="unsigned_summary_file",
            folder=f"file_summaries/{uid}",
            resource_type="auto"
        )
        file_url = upload_result["secure_url"]
        file_name = upload_result["original_filename"]

        # --- 3. Rewind File (Crucial!) ---
        file.seek(0)

        # --- 4. Extract Text OR Flag for Visual Mode ---
        extracted_text = ""
        is_visual_file = False 

        if filename.endswith(".pdf"):
            extracted_text += extract_text_from_pdf(file_url)
            # If text extraction fails, treat as a scanned PDF (Visual)
            if not extracted_text.strip():
                is_visual_file = True
        
        elif filename.endswith(".docx"):
            extracted_text += extract_text_from_docx(file_url)
        elif filename.endswith(".pptx"):
            extracted_text += extract_text_from_pptx(file_url)
        elif filename.endswith((".xls", ".xlsx")):
            extracted_text += extract_text_from_excel(file_url)
        
        # Handle Standard Images
        elif filename.endswith((".jpg", ".jpeg", ".png", ".webp", ".heic")):
            is_visual_file = True
        
        else:
            return jsonify({"error": "Unsupported file type"}), 400


        # --- 5. Generate Summary (With 3-Layer Backup) ---
        final_answer = ""

        # === METHOD A: TEXT-BASED SUMMARY ===
        if extracted_text.strip() and not is_visual_file:
            # 1. Chunking
            def chunk_text(text, max_words=1000): # Increased chunk size slightly
                words = text.split()
                return [' '.join(words[i:i + max_words]) for i in range(0, len(words), max_words)]

            chunks = chunk_text(extracted_text)
            summaries = []
            
            # 2. Summarize Chunks (Primary Model Only)
            for chunk in chunks:
                chunk_prompt = f"Summarize this section clearly:\n\n{chunk.strip()}"
                try:
                    response = model_pdf.generate_content(chunk_prompt)
                    summaries.append(response.text.strip())
                except Exception as e:
                    summaries.append("") # Skip bad chunks or handle differently

            # 3. Final Compilation
            full_summary_prompt = (
                "You are Felix. Create a detailed academic summary from these notes:\n\n" +
                "\n\n".join(summaries)
            )

            try:
                # LAYER 1: Primary Model
                final_response = model.generate_content(full_summary_prompt)
                raw_answer = final_response.text.strip()
            
            except Exception as e:
                print(f"Primary Text Summary Failed: {e}")
                try:
                    # LAYER 2: First Backup
                    raw_answer = generate_backup_response(full_summary_prompt)
                except Exception as e2:
                    print(f"First Backup Failed: {e2}")
                    # LAYER 3: Final 6-Key Backup (Last Resort)
                    raw_answer = final_backup_response(full_summary_prompt)

            final_answer = clean_response(raw_answer)

        # === METHOD B: VISUAL SUMMARY (Images/Scanned PDFs) ===
        elif is_visual_file:
            try:
                # Prepare visual data
                file_data = file.read()
                visual_prompt = "You are Felix. Analyze this image/document. Provide a detailed academic summary."
                
                # Input for Gemini (List of [Prompt, ImageDict])
                gemini_input = [visual_prompt, {"mime_type": mime_type, "data": file_data}]

                # LAYER 1: Primary Model
                response = model_pdf.generate_content(gemini_input)
                raw_answer = response.text.strip()

            except Exception as e:
                print(f"Primary Visual Summary Failed: {e}")
                # Prepare input again (in case it was modified)
                file.seek(0)
                file_data = file.read()
                gemini_input = [visual_prompt, {"mime_type": mime_type, "data": file_data}]

                try:
                    # LAYER 2: First Backup
                    raw_answer = generate_backup_response(gemini_input)
                except Exception as e2:
                    print(f"First Visual Backup Failed: {e2}")
                    # LAYER 3: Final 6-Key Backup (Last Resort)
                    raw_answer = final_backup_response(gemini_input)
            
            final_answer = clean_response(raw_answer)

        else:
            return jsonify({"error": "Could not extract text from this file."}), 400


        # --- 6. Save to Firestore ---
        db.collection("users").document(uid).collection("summaries").add({
            "summary": final_answer,
            "created_at": datetime.utcnow().isoformat(),
            "source_file": file_name,
            "file_url": file_url
        })

        return jsonify({"response": final_answer, "file_url": file_url}), 200

    except Exception as e:
        print("Error in /media_summary:", str(e))
        return jsonify({"error": str(e)}), 500

@app.route('/get_summaries', methods=['GET'])
def get_summaries():
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        db = firestore.client()
        summaries_ref = db.collection("users").document(uid).collection("summaries").order_by("created_at", direction=firestore.Query.DESCENDING)
        summaries_docs = summaries_ref.stream()

        summaries = []
        for doc in summaries_docs:
            data = doc.to_dict()
            summaries.append({
                "id": doc.id,
                "summary": data.get("summary", ""),
                "created_at": data.get("created_at", "")
            })

        return jsonify({"summaries": summaries})

    except Exception as e:
        return jsonify({"error": str(e)}), 500
@app.route('/delete_summaries/<summary_id>', methods=['DELETE'])
def delete_summary(summary_id):
    try:
        decoded = verify_token_from_request()
        uid = decoded['uid']
        doc_ref = db.collection("users").document(uid).collection("summaries").document(summary_id)
        doc = doc_ref.get()
        if not doc.exists:
            return jsonify({"error": "Quiz set not found"}), 404
        doc_ref.delete()
        return jsonify({"ok": True, "message": "summary successfully deleted"}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500



@app.route('/chat_quiz', methods=['POST'])
def create_chat_quiz():
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]
        update_streak(uid)

        content = request.json.get("text", "")
        if not content:
            return jsonify({"error": "No input text provided"}), 400

        # --- 1. Stronger Prompt ---
        prompt = f"""
        You are ACE, a helpful AI tutor. Based on the content below, generate 10 quiz flashcards.
        Strictly follow this JSON schema:
        [
          {{
            "question": "string",
            "options": {{ "A": "...", "B": "...", "C": "...", "D": "..." }},
            "answer": "A"
          }}
        ]
        Do not add any Markdown formatting like ```json. Just return the raw JSON array.
        
        Content:
        {content}
        """

        quiz_data = []

        try:
            # --- 2. Primary Generation ---
            response = model_cards.generate_content(prompt)
            quiz_raw = response.text.strip()
            
            # Robust Parsing
            start_index = quiz_raw.find('[')
            end_index = quiz_raw.rfind(']')
            
            if start_index != -1 and end_index != -1:
                json_str = quiz_raw[start_index : end_index + 1]
                quiz_data = json.loads(json_str)
            else:
                raise ValueError("No JSON list found in response")

        except Exception as e:
            print(f"Primary Gemini failed: {e}")
            
            # --- 3. Backup Strategy (Added) ---
            try:
                # Call your global backup function
                backup_raw = final_backup_response(prompt)

                # Parse Backup Response
                start_index = backup_raw.find('[')
                end_index = backup_raw.rfind(']')
                
                if start_index != -1 and end_index != -1:
                    json_str = backup_raw[start_index : end_index + 1]
                    quiz_data = json.loads(json_str)
                else:
                    print("Backup returned invalid JSON")
            except Exception as backup_error:
                print(f"Backup also failed: {backup_error}")
                pass

        # --- 4. Validation & Saving ---
        # Ensure we have data (from either Primary or Backup)
        if not isinstance(quiz_data, list):
            quiz_data = []

        if not quiz_data:
            return jsonify({"error": "Failed to generate valid quiz data from AI."}), 500

        # 5. Save structured quiz to Firestore
        doc_ref = db.collection("users").document(uid).collection("quiz").document()
        doc_ref.set({
            "quiz": quiz_data,
            "created_at": datetime.utcnow().isoformat()
        })

        # 6. Return result
        return jsonify({"quiz": quiz_data}), 200

    except Exception as e:
        return jsonify({"error": str(e)}), 500



@app.route('/media_quiz', methods=['POST'])
def media_quiz():
    try:
        # --- 1. Auth & Setup ---
        decoded = verify_token_from_request()
        uid = decoded["uid"]
        update_streak(uid)

        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400

        file = request.files['file']
        filename = file.filename.lower()

        # --- 2. Upload to Cloudinary ---
        upload_result = cloudinary.uploader.upload(
            file,
            upload_preset="unsigned_quiz_uploads",
            folder=f"file_quizzes/{uid}",
            resource_type="auto"
        )

        file_url = upload_result["secure_url"]
        file_name = upload_result["original_filename"]

        # --- 3. Smart Extraction (Text or Image) ---
        extracted_text = ""
        image_part = None  # Stores the image object if we are in Image Mode

        try:
            if filename.endswith(".pdf"):
                extracted_text = extract_text_from_pdf(file_url)
            elif filename.endswith(".docx"):
                extracted_text = extract_text_from_docx(file_url)
            elif filename.endswith(".pptx"):
                extracted_text = extract_text_from_pptx(file_url)
            elif filename.endswith((".xls", ".xlsx")):
                extracted_text = extract_text_from_excel(file_url)
            
            # ✅ Image Support (Native Vision)
            elif filename.endswith((".png", ".jpg", ".jpeg", ".webp")):
                response = requests.get(file_url)
                image_part = Image.open(BytesIO(response.content))
                extracted_text = "IMAGE_MODE" 
            else:
                return jsonify({"error": "Unsupported file type"}), 400
        except Exception as e:
            print(f"Extraction error: {e}")
            return jsonify({"error": "Failed to read document content"}), 400

        if not extracted_text:
            return jsonify({"error": "No readable content found"}), 400

        # --- 4. AI Generation ---
        base_prompt = (
            "You are ACE, a helpful AI tutor. Based on the study material provided, "
            "generate exactly 10 multiple-choice quiz flashcards in JSON format.\n"
            "Strictly follow this schema:\n"
            "[{\"question\": \"...\", \"options\": {\"A\": \"...\", \"B\": \"...\", \"C\": \"...\", \"D\": \"...\"}, \"answer\": \"A\"}]\n"
            "Return ONLY raw JSON. No Markdown."
        )

        quiz_data = []

        try:
            # --- Primary Attempt ---
            if image_part:
                # Image Mode (Send Prompt + Image)
                response = model_cards.generate_content([base_prompt, image_part])
            else:
                # Text Mode (Limited to 7500 chars)
                full_prompt = f"{base_prompt}\n\nText:\n{extracted_text[:7500]}"
                response = model_cards.generate_content(full_prompt)

            raw_response = response.text.strip()
            
            # Robust Parsing
            start = raw_response.find('[')
            end = raw_response.rfind(']')
            if start != -1 and end != -1:
                quiz_data = json.loads(raw_response[start : end + 1])
            else:
                raise ValueError("No JSON list found")

        except Exception as e:
            print(f"Gemini quiz generation failed: {e}")
            
            # --- Backup Strategy (Using KeyManager) ---
            try:
                # 1. Determine Input for Backup
                if image_part:
                    backup_input = [base_prompt, image_part]
                else:
                    backup_input = f"{base_prompt}\n\nText:\n{extracted_text[:7500]}"
                
                # 2. Call the NEW global backup function
                backup_raw = final_backup_response(backup_input) # <--- RENAMED CALL
                
                # 3. Parse Backup Response
                start = backup_raw.find('[')
                end = backup_raw.rfind(']')
                if start != -1 and end != -1:
                    quiz_data = json.loads(backup_raw[start : end + 1])
                else:
                    print("Backup returned invalid JSON")
            except Exception as backup_error:
                print(f"Backup also failed: {backup_error}")
                pass

        # --- 6. Validation ---
        if not isinstance(quiz_data, list) or len(quiz_data) == 0:
             return jsonify({"error": "Could not generate valid quiz."}), 500

        # --- 7. Save & Return ---
        db.collection("users").document(uid).collection("quiz").add({
            "quiz": quiz_data,
            "created_at": datetime.utcnow().isoformat(),
            "source_file": file_name,
            "file_url": file_url
        })

        return jsonify({"quiz": quiz_data, "file_url": file_url}), 200

    except Exception as e:
        print("Error in /media_quiz:", str(e))
        return jsonify({"error": str(e)}), 500

# ---------- Get All Quizzes ----------
@app.route('/get_chat_quiz', methods=['GET'])
def get_chat_quiz():
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        quiz_docs = db.collection("users").document(uid).collection("quiz").order_by("created_at").get()

        quiz_list = []
        for doc in quiz_docs:
            data = doc.to_dict()
            quiz_list.append({
                "id": doc.id,
                "quiz": data.get("quiz"),
                "created_at": data.get("created_at")
            })

        return jsonify({"quiz": quiz_list}), 200

    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ---------- Get One Quiz by ID ----------
@app.route('/get_chat_quiz/<quiz_id>', methods=['GET'])
def get_single_quiz(quiz_id):
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]
        update_streak(uid)

        doc_ref = db.collection("users").document(uid).collection("quiz").document(quiz_id)
        doc = doc_ref.get()

        if not doc.exists:
            return jsonify({"error": "Quiz not found"}), 404

        quiz_data = doc.to_dict()
        return jsonify({
            "id": doc.id,
            "quiz": quiz_data.get("quiz"),
            "created_at": quiz_data.get("created_at")
        }), 200

    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ---------- Delete Quiz ----------
@app.route('/delete_quiz/<quiz_id>', methods=['DELETE'])
def delete_quiz(quiz_id):
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        doc_ref = db.collection("users").document(uid).collection("quiz").document(quiz_id)
        doc = doc_ref.get()

        if not doc.exists:
            return jsonify({"error": "Quiz set not found"}), 404

        doc_ref.delete()
        return jsonify({"ok": True, "message": "Quiz successfully deleted"}), 200

    except Exception as e:
        return jsonify({"error": str(e)}), 500
    

@app.route('/save_quiz_score', methods=['POST'])
def save_quiz_score():
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        data = request.json
        quiz_id = data.get("quiz_id")  # Could be PDF ID or timestamp
        score = data.get("score")
        total = data.get("total")

        # Reference quiz scores collection
        score_ref = db.collection("users").document(uid).collection("quiz_scores").document(quiz_id)

        # Update if exists, else create new
        score_ref.set({
            "score": score,
            "total": total,
            "percentage": (score / total) * 100 if total > 0 else 0,
            "updated_at": datetime.utcnow()
        }, merge=True)

        return jsonify({"message": "Score saved successfully"}), 200

    except Exception as e:
        return jsonify({"error": str(e)}), 500





@app.route('/media_flashcards', methods=['POST'])
def media_flashcards():
    try:
        # --- 1. Auth & Setup ---
        decoded = verify_token_from_request()
        uid = decoded["uid"]
        update_streak(uid)

        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400

        file = request.files['file']
        filename = file.filename.lower()

        # --- 2. Extract text FIRST (before uploading) ---
        extracted_text = ""
        image_part = None

        try:
            if filename.endswith(".pdf"):
                # Read directly from upload
                file.seek(0)
                reader = PdfReader(file)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        extracted_text += page_text + "\n"
                
            elif filename.endswith(".docx"):
                file.seek(0)
                doc = Document(file)
                extracted_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                
            elif filename.endswith(".pptx"):
                file.seek(0)
                prs = Presentation(file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
                            
            elif filename.endswith((".xls", ".xlsx")):
                file.seek(0)
                wb = load_workbook(file)
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " ".join(str(cell) for cell in row if cell)
                        if row_text.strip():
                            extracted_text += row_text + "\n"
            
            elif filename.endswith((".png", ".jpg", ".jpeg", ".webp")):
                file.seek(0)
                image_part = Image.open(file)
                extracted_text = "IMAGE_MODE"
            else:
                return jsonify({"error": "Unsupported file type"}), 400
                
        except Exception as e:
            print(f"Extraction error: {e}")
            return jsonify({"error": "Failed to read document text"}), 400

        if not extracted_text:
            return jsonify({"error": "No readable content found"}), 400

        # --- 3. NOW upload to Cloudinary for storage ---
        file.seek(0)  # Reset file pointer
        upload_result = cloudinary.uploader.upload(
            file,
            upload_preset="unsigned_flashcards",
            folder=f"flashcards/{uid}",
            resource_type="auto"
        )

        file_url = upload_result["secure_url"]
        file_name = upload_result["original_filename"]

        # --- 4. AI Generation ---
        base_prompt = (
            "You are ACE, a helpful academic AI assistant. "
            "Generate exactly 10 high-quality flashcards based on the study material. "
            "Return ONLY a valid JSON array of objects with 'question' and 'answer' fields. "
            "Do not include Markdown formatting (no ```json)."
        )

        flashcards = []

        try:
            # --- Primary Attempt ---
            if image_part:
                # Image Mode
                response = model_pdf.generate_content([base_prompt, image_part])
            else:
                # Text Mode (Limited to 7500 chars)
                final_prompt = f"{base_prompt}\n\nText:\n{extracted_text[:7500]}"
                response = model_pdf.generate_content(final_prompt)

            raw_response = response.text.strip()
            
            # Parsing
            start = raw_response.find('[')
            end = raw_response.rfind(']')
            if start != -1 and end != -1:
                flashcards = json.loads(raw_response[start : end + 1])
            else:
                raise ValueError("No JSON list found")

        except Exception as e:
            print(f"Primary generation failed: {e}")
            
            # --- Backup Strategy ---
            try:
                # 1. Determine Input for Backup
                if image_part:
                    # Multimodal Backup
                    backup_input = [base_prompt, image_part]
                else:
                    # Text Backup
                    backup_input = f"{base_prompt}\n\nText:\n{extracted_text[:7500]}"
                
                # 2. Call the global backup function
                backup_raw = final_backup_response(backup_input)

                # 3. Parse Backup Response
                start = backup_raw.find('[')
                end = backup_raw.rfind(']')
                if start != -1 and end != -1:
                    flashcards = json.loads(backup_raw[start : end + 1])
                else:
                    print("Backup returned invalid JSON")
            except Exception as backup_error:
                print(f"Backup also failed: {backup_error}")
                pass 

        # --- 5. Validation ---
        if not isinstance(flashcards, list):
            flashcards = []
        
        if len(flashcards) > 15: 
            flashcards = flashcards[:15]

        if not flashcards:
            return jsonify({"error": "AI could not generate flashcards from this file"}), 500

        # --- 6. Save & Return ---
        db.collection("users").document(uid).collection("flashcards").add({
            "flashcards": flashcards,
            "created_at": datetime.utcnow().isoformat(),
            "source_file": file_name,
            "file_url": file_url
        })

        return jsonify({"flashcards": flashcards, "file_url": file_url}), 200

    except Exception as e:
        print(f"Media Route Error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/chat_flashcards', methods=['POST'])
def chat_flashcards():
    try:
        decoded = verify_token_from_request()
        uid = decoded['uid']
        update_streak(uid)
        
        data = request.get_json().get("text", "").strip()
        if not data:
            return jsonify({"error": "no text provided"}), 400

        flashcard_prompt = (
            "You are ACE, a helpful academic AI assistant. "
            "Based on the content below, generate 10 flashcards in JSON format. "
            "Return ONLY a valid JSON array of objects, where each object has 'question' and 'answer'. "
            "Do not include Q:/A: labels, explanations, or extra text. "
            "For example:\n"
            "[{\"question\": \"What is X?\", \"answer\": \"X is ...\"}, ...]\n\n"
            f"Text:\n\n{data}"
        )

        flashcards = [] # Initialize variable

        # --- 1. Generation & Parsing ---
        try:
            response = model_cards.generate_content(flashcard_prompt)
            flashcard_raw = response.text.strip()
            
            # Robust Extraction: Find the list brackets directly
            start = flashcard_raw.find('[')
            end = flashcard_raw.rfind(']')
            
            if start != -1 and end != -1:
                json_str = flashcard_raw[start : end + 1]
                flashcards = json.loads(json_str)
            else:
                raise ValueError("No JSON list found in response")

        except Exception as e:
            print(f"Primary Flashcard generation failed: {e}")
            
            # --- 2. Backup Handling ---
            # Assuming generate_backup_response returns a string
            backup_raw = final_backup_response(flashcard_prompt)
            
            # Try to parse backup, or fail gracefully
            try:
                start = backup_raw.find('[')
                end = backup_raw.rfind(']')
                if start != -1 and end != -1:
                    flashcards = json.loads(backup_raw[start : end + 1])
                else:
                    # Fallback if backup is not JSON: wrap raw text in a single card
                    flashcards = [{"question": "Summary", "answer": backup_raw}]
            except:
                flashcards = [{"question": "Error", "answer": "Could not generate flashcards."}]

        # --- 3. Save & Return (NOW UN-INDENTED) ---
        # This code now runs regardless of whether the Try or Except block executed.
        
        # Validation: Ensure we actually have a list before saving
        if not isinstance(flashcards, list):
            flashcards = []

        db.collection("users").document(uid).collection("flashcards").add({
            "flashcards": flashcards,
            "created_at": datetime.utcnow().isoformat()
        })

        return jsonify({"flashcards": flashcards}), 200

    except Exception as e:
        print(f"Critical Route Error: {e}")
        return jsonify({"error": str(e)}), 500


    

@app.route('/get_flashcards', methods=['GET'])
def get_flashcards():
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]
        update_streak(uid)

        flashcards = (
            db.collection("users")
              .document(uid)
              .collection("flashcards")
              .order_by("created_at")
              .get()
        )

        flashcard_sets = []
        for doc in flashcards:
            data = doc.to_dict()
            flashcard_sets.append({
                "id": doc.id,
                "flashcards": data.get("flashcards", []),
                "created_at": data.get("created_at")
            })

        return jsonify({"flashcards": flashcard_sets}), 200

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/update_flashcards/<flashcard_id>', methods=['PUT'])
def update_flashcards(flashcard_id):
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        data = request.get_json() or {}
        new_flashcards = data.get("flashcards")

        if not new_flashcards or not isinstance(new_flashcards, list):
            return jsonify({"error": "Flashcards must be provided as a list"}), 400

        doc_ref = db.collection("users").document(uid).collection("flashcards").document(flashcard_id)
        doc = doc_ref.get()

        if not doc.exists:
            return jsonify({"error": "Flashcard set not found"}), 404

        doc_ref.update({
            "flashcards": new_flashcards,
            "updated_at": firestore.SERVER_TIMESTAMP
        })

        return jsonify({"ok": True, "message": "Flashcards updated"}), 200

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/flashcards/<flashcard_id>', methods=['DELETE'])
def delete_flashcards(flashcard_id):
    try:
        decoded = verify_token_from_request()
        uid = decoded['uid']
        flashcard_ref = db.collection("users").document(uid).collection("flashcards").document(flashcard_id)
        flashcard_ref.delete()
        return jsonify({ 'ok': True, "message": "flashcard deleted sucessfully"})
    except Exception as e :
        return jsonify ({'error': str(e)}),400



# ---------- Create Study Reminder ----------
@app.route("/reminders/study-plan", methods=["POST"])
def create_study_plan():
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]
        update_streak(uid)

        # Accept form-data
        title = request.form.get("title")
        start_date = request.form.get("start_date")
        due_date = request.form.get("due_date")
        description = request.form.get("description", "")
        uploaded_files = request.files.getlist("files")  # list of files

        file_links = []
        extracted_text = ""

        # ---------- Upload files to Cloudinary ----------
        for file in uploaded_files:
            upload_result = cloudinary.uploader.upload(
                file,
                upload_preset="unsigned_study_plans",
                folder=f"study_plans/{uid}",
                resource_type="auto"
            )
            file_url = upload_result["secure_url"]
            file_name = upload_result["original_filename"]

            file_links.append({"name": file_name, "url": file_url})

            # ---------- Extract text (PDF/Docs only) ----------
            # Note: We removed extract_text_from_image to avoid Tesseract crashes.
            # If you need image support here later, we can add Single-Image Native Vision.
            try:
                if file_name.lower().endswith(".pdf"):
                    extracted_text += extract_text_from_pdf(file_url)
                elif file_name.lower().endswith((".docx", ".doc")):
                    extracted_text += extract_text_from_docx(file_url)
            except Exception as e:
                print(f"Failed to extract text from {file_name}: {e}")

        # ---------- Generate study plan ----------
        prompt = f"""
        You are a helpful AI tutor.
        I have an upcoming test: {title}.
        Start date: {start_date}.
        Due date: {due_date}.
        Description: {description}.
        Notes/Materials: {extracted_text[:7500]} 

        Please create a day-by-day study plan starting from the Start date in this exact format:
        Day 1: Topic 1, Topic 2
        Day 2: Topic 3, Topic 4
        Day 3: Topic 5
        ...
        Continue until the Due date.
        """

        try:
            # Primary Generation
            response = reminder_model.generate_content(prompt)
            text_output = response.text
        except Exception as e:
            print(f"Primary Study Plan Gen failed: {e}")
            # Final Backup Logic
            text_output = final_backup_response(prompt)

        # Parse the result (whether from Primary or Backup)
        study_plan = parse_study_plan(text_output)

        # ---------- Save to Firestore ----------
        reminder_ref = db.collection("users").document(uid).collection("reminders").document()
        reminder_ref.set({
            "title": title,
            "due_date": due_date,
            "description": description,
            "files": file_links,
            "study_plan": study_plan,
            "completed": False,
            "created_at": firestore.SERVER_TIMESTAMP,
        })

        return jsonify({"ok": True, "study_plan": study_plan, "files": file_links})

    except Exception as e:
        print(f"Study Plan Error: {e}")
        return jsonify({"error": str(e)}), 400



# ---------- Get Reminders ----------
@app.route("/get_reminders", methods=["GET"])
def get_reminders():
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        reminders_ref = db.collection("users").document(uid).collection("reminders")
        docs = reminders_ref.order_by("created_at", direction=firestore.Query.DESCENDING).stream()

        reminders = []
        for doc in docs:
            data = doc.to_dict()
            data["id"] = doc.id
            reminders.append(data)

        return jsonify({"ok": True, "reminders": reminders})

    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 400



# ---------- Complete Reminder ----------
@app.route("/reminders/<reminder_id>/complete", methods=["POST"])
def complete_reminder(reminder_id):
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        reminder_ref = db.collection("users").document(uid).collection("reminders").document(reminder_id)
        reminder_ref.update({"completed": True})

        return jsonify({"ok": True, "message": "Reminder marked as completed"})
    except Exception as e:
        return jsonify({"error": str(e)}), 400


# ---------- Delete Reminder ----------
@app.route("/reminders/<reminder_id>", methods=["DELETE"])
def delete_reminder(reminder_id):
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        reminder_ref = db.collection("users").document(uid).collection("reminders").document(reminder_id)
        reminder_doc = reminder_ref.get()

        if reminder_doc.exists:
            reminder_data = reminder_doc.to_dict()

            # Optional: delete files from Cloudinary
            if "files" in reminder_data and isinstance(reminder_data["files"], list):
                for f in reminder_data["files"]:
                    try:
                        file_url = f.get("url")
                        if file_url:
                            # Extract public_id safely (assuming format: https://res.cloudinary.com/<cloud_name>/.../study_plans/<uid>/<public_id>.<ext>)
                            public_id = file_url.split("/")[-1].split(".")[0]
                            cloudinary.uploader.destroy(f"study_plans/{uid}/{public_id}")
                    except Exception as inner_e:
                        print(f"Cloudinary delete error: {inner_e}")

            # Delete Firestore doc
            reminder_ref.delete()

            return jsonify({"ok": True, "message": "Reminder and files deleted successfully"})
        else:
            return jsonify({"error": "Reminder not found"}), 404

    except Exception as e:
        return jsonify({"error": str(e)}), 400

    
#---Content-Recommender-- 
@app.route("/recommendations", methods=["GET"])
def recommend_content():
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        # 1. Get user profile
        user_doc = db.collection("users").document(uid).get()
        if not user_doc.exists:
            return jsonify({"error": "User profile not found"}), 404
        user_data = user_doc.to_dict()

        course = user_data.get("course_of_study", "")
        subjects = user_data.get("subject", [])
        if isinstance(subjects, str):
            subjects = [subjects]

        if not course and not subjects:
            return jsonify({"error": "No course_of_study or subjects found in profile"}), 400

        # 2. Build query
        query_topics = [course] + subjects
        search_query = " ".join(query_topics)

        # 3. Fetch YouTube videos
        youtube_url = "https://www.googleapis.com/youtube/v3/search"
        params = {
            "part": "snippet",
            "q": search_query,
            "key": YOUTUBE_API_KEY,
            "maxResults": 5,
            "type": "video"
        }
        yt_response = requests.get(youtube_url, params=params).json()

        if "items" not in yt_response:
            return jsonify({"error": "No results from YouTube"}), 400

        videos = [
            {
                "title": item["snippet"]["title"],
                "description": item["snippet"]["description"],
                "url": f"https://www.youtube.com/watch?v={item['id']['videoId']}"
            }
            for item in yt_response["items"]
        ]

        # 4. Ask Gemini for JSON output
        prompt = f"""
        The student is studying: {course}.
        Their subjects are: {', '.join(subjects)}.
        Here are some YouTube videos: {json.dumps(videos, indent=2)}.

        Please return the top 3 recommended videos as a STRICT JSON list.
        Each entry should have:
        - "title"
        - "url"
        - "reason"
        """

        try:
            response = model.generate_content(prompt)
        except Exception as e:
            print(f"Gemini failed oo {e}")
            response = rec_backup_response(prompt)
        
        recommendations = []
        try:
            recommendations = json.loads(response.text)
        except Exception:
            recommendations = videos[:3]  # fallback to top 3 raw videos

        # 5. Save recommendations to Firestore
        recs_ref = db.collection("users").document(uid).collection("recommendations")
        rec_doc = {
            "created_at": datetime.utcnow().isoformat(),
            "course": course,
            "subjects": subjects,
            "recommendations": recommendations
        }
        recs_ref.add(rec_doc)

        return jsonify({
            "ok": True,
            "recommendations": recommendations,
            "raw_videos": videos
        })

    except Exception as e:
        return jsonify({"error": str(e)}), 400




@app.route("/predict_performance", methods=["GET"])
def predict_performance():
    try:
        # Verify user authentication
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        # 1️Fetch user data
        user_ref = db.collection("users").document(uid)
        user_doc = user_ref.get()
        if not user_doc.exists:
            return jsonify({"error": "User not found"}), 404

        user_data = user_doc.to_dict()

        # Compute age from date_of_birth (expecting 'YYYY-MM-DD' or ISO format)
        default_age = 17
        dob_str = user_data.get("date_of_birth")  # your field
        age = default_age
        if dob_str:
            try:
                # handles "YYYY-MM-DD" and ISO datetime strings
                dob = datetime.fromisoformat(dob_str).date()
            except Exception:
                try:
                    dob = datetime.strptime(dob_str, "%Y-%m-%d").date()
                except Exception:
                    dob = None
            if dob:
                today = datetime.utcnow().date()
                age = today.year - dob.year - ((today.month, today.day) < (dob.month, dob.day))
                if age < 0:
                    # safeguard if DOB is in future by mistake
                    age = 0

        # Prepare ML model input
        gender_val = user_data.get("gender", "")
        gender_val = gender_val.lower() if isinstance(gender_val, str) else ""
        ml_input = {
            "study_hours_per_week": round(user_data.get("study_hours_per_week", 0) / 40, 3),
            "sleep_hours_per_day": round(user_data.get("sleep_hours_per_day", 8) / 10, 3),
            "attendance_percentage": round(user_data.get("attendance_percentage", 0) / 100, 3),
            "assignments_completed": round(user_data.get("assignment_completed", 0), 3),
            "participation_level": user_data.get("participation_level", 0),
            "Age": int(age),
            "Gender": 1 if gender_val == "male" else 0,
            "StudyTimeWeekly": round(user_data.get("study_hours_per_week", 0), 2),
            "Absences": user_data.get("Absences", 0),
            "Tutoring": user_data.get("Tutoring", 0)
        }

        df = pd.DataFrame([ml_input])

        # Run prediction
        prediction = pipeline.predict(df)
        proba = pipeline.predict_proba(df)[0]
        confidence = float(proba[1] * 100)
        fail_confidence = float((1 - proba[1]) * 100)

        # Generate message
        if int(prediction[0]) == 1:
            message = f"Our model predicts you'll PASS your next exam 🎯 (Confidence: {confidence:.2f}%). Keep it up!"
        else:
            message = f"The model predicts you might fail your next exam 📊 (Confidence: {fail_confidence:.2f}%). Let's improve your learning habits."

            reasons = []
            if ml_input["study_hours_per_week"] < 0.4:
                reasons.append("increase your weekly study hours")
            if ml_input["attendance_percentage"] < 0.7:
                reasons.append("improve your class attendance")
            if ml_input["assignments_completed"] < 0.5:
                reasons.append("complete more assignments/quizzes on time")
            if ml_input["sleep_hours_per_day"] < 0.5:
                reasons.append("maintain a healthier sleep schedule")

            if reasons:
                message += " To improve, try to " + ", ".join(reasons) + "."


        def to_native(value):
            if isinstance(value, (np.floating, np.float32, np.float64)):
                return float(value)
            elif isinstance(value, (np.integer, np.int32, np.int64)):
                return int(value)
            elif isinstance(value, dict):
                return {k: to_native(v) for k, v in value.items()}
            elif isinstance(value, list):
                return [to_native(v) for v in value]
            else:
                return value


        prediction_data = {
            "timestamp": datetime.utcnow().isoformat(),
            "prediction": int(prediction[0]),
            "confidence": float(confidence),
            "message": message,
            "input_features": to_native(ml_input)
        }

        user_ref.collection("predictions").add(prediction_data)

        return jsonify({
            "ok": True,
            "prediction": int(prediction[0]),
            "confidence": confidence,
            "message": message,
            "ml_input": to_native(ml_input)
        }), 200

    except Exception as e:
        return jsonify({"error": str(e)}), 500



@app.route("/gamification/<uid>", methods=["GET", "POST"])
def gamification(uid):
    try:

        decoded = verify_token_from_request()
        if decoded["uid"] != uid:
            return jsonify({"error": "Unauthorized"}), 403

        user_ref = db.collection("users").document(uid)
        gamification_ref = user_ref.collection("meta").document("gamification")

        # Frontend sends badge or level updates
        if request.method == "POST":
            data = request.get_json() or {}
            add_badge = data.get("badge")
            new_level = data.get("level")

            # Fetch current gamification data
            doc = gamification_ref.get()
            gamification_data = doc.to_dict() if doc.exists else {
                "level": 1,
                "badges": []
            }

            # Frontend decides when to give badges
            if add_badge and add_badge not in gamification_data["badges"]:
                gamification_data["badges"].append(add_badge)

            # Optional manual level update
            if new_level and isinstance(new_level, int):
                gamification_data["level"] = new_level

            gamification_data["last_activity"] = firestore.SERVER_TIMESTAMP
            gamification_ref.set(gamification_data)

            return jsonify({"ok": True, "gamification": gamification_data}), 200

        # GET → Retrieve gamification info + streak from user profile
        elif request.method == "GET":
            doc = gamification_ref.get()
            gamification_data = doc.to_dict() if doc.exists else {
                "level": 1,
                "badges": []
            }

            # Add streak info from user profile
            user_doc = user_ref.get()
            streak_info = {"streak": 0, "longest_streak": 0}
            if user_doc.exists:
                user_data = user_doc.to_dict()
                streak_info["streak"] = user_data.get("streak", 0)
                streak_info["longest_streak"] = user_data.get("longest_streak", 0)

            gamification_data.update(streak_info)

            return jsonify({"gamification": gamification_data}), 200

    except Exception as e:
        print("Gamification error:", e)
        return jsonify({"error": str(e)}), 500


@app.route("/library/<uid>", methods=["GET", "POST", "DELETE"])
def library(uid):
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]
        user_library_ref = db.collection("users").document(uid).collection("library")

        # -----------------------
        # POST → Upload documents
        # -----------------------
        if request.method == "POST":
            title = request.form.get("title")
            uploaded_files = request.files.getlist("files")

            if not uploaded_files:
                return jsonify({"error": "No files uploaded"}), 400

            file_links = []
            for file in uploaded_files:
                upload_result = cloudinary.uploader.upload(
                    file,
                    upload_preset="unsigned_library",
                    folder=f"library/{uid}",
                    resource_type="auto",  # handles pdf, image, docx, etc.
                )
                file_links.append({
                    "name": upload_result.get("original_filename"),
                    "url": upload_result.get("secure_url"),
                    "type": upload_result.get("resource_type"),
                })

            doc_ref = user_library_ref.document()
            doc_ref.set({
                "title": title,
                "files": file_links,
                "created_at": firestore.SERVER_TIMESTAMP,
            })

            return jsonify({"ok": True, "message": "Files uploaded successfully", "file_links": file_links}), 200

        # -----------------------
        # GET → Retrieve all documents
        # -----------------------
        elif request.method == "GET":
            docs = user_library_ref.order_by("created_at", direction=firestore.Query.DESCENDING).stream()
            library_items = [{"id": doc.id, **doc.to_dict()} for doc in docs]
            return jsonify({"ok": True, "library": library_items}), 200

        # -----------------------
        # DELETE → Delete all library entries
        # -----------------------
        elif request.method == "DELETE":
            docs = user_library_ref.stream()
            for doc in docs:
                doc.reference.delete()
            return jsonify({"ok": True, "message": "All library documents deleted"}), 200

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# -----------------------
# PUT → Update specific library document (title or files)
# -----------------------
@app.route("/library_update/<uid>/<doc_id>", methods=["PUT"])
def update_library(uid, doc_id):
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]
        title = request.form.get("title")
        uploaded_files = request.files.getlist("files")

        doc_ref = db.collection("users").document(uid).collection("library").document(doc_id)

        update_data = {}
        if title:
            update_data["title"] = title
        if uploaded_files:
            file_links = []
            for file in uploaded_files:
                upload_result = cloudinary.uploader.upload(
                    file,
                    upload_preset="unsigned_library",
                    folder=f"library/{uid}",
                    resource_type="auto",
                )
                file_links.append({
                    "name": upload_result.get("original_filename"),
                    "url": upload_result.get("secure_url"),
                    "type": upload_result.get("resource_type"),
                })
            update_data["files"] = firestore.ArrayUnion(file_links)

        if update_data:
            doc_ref.update(update_data)
            return jsonify({"ok": True, "message": "Library entry updated successfully"}), 200
        else:
            return jsonify({"error": "No updates provided"}), 400

    except Exception as e:
        return jsonify({"error": str(e)}), 400


# -----------------------
# DELETE → Delete specific library document
# -----------------------
@app.route("/library_delete/<uid>/<doc_id>", methods=["DELETE"])
def delete_library_item(uid, doc_id):
    try:
        decoded = verify_token_from_request()
        uid = decoded["uid"]

        doc_ref = db.collection("users").document(uid).collection("library").document(doc_id)
        doc_ref.delete()

        return jsonify({"ok": True, "message": "Library document deleted successfully"}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 400



# ---------- Run ----------
if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)
